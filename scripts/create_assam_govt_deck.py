from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "khozo-assam-govt-deck.pptx"

GREEN = RGBColor(103, 164, 38)
DARK_GREEN = RGBColor(78, 128, 25)
LIGHT_GREEN = RGBColor(234, 244, 221)
INK = RGBColor(28, 37, 48)
MUTED = RGBColor(86, 99, 115)
SLATE = RGBColor(226, 232, 240)
BLUE = RGBColor(37, 99, 235)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)

SRC_PDF = "Source: Khozo (1).pdf"
SRC_DOCX = "Source: khozo_proposal_v1.0.docx"


def textbox(slide, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def fill_rect(slide, x, y, w, h, color, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def header(slide, title, section, page):
    textbox(slide, 0.62, 0.32, 3.8, 0.25, section.upper(), 8.5, GREEN, True)
    textbox(slide, 0.58, 0.58, 9.2, 0.55, title, 24, INK, True)
    fill_rect(slide, 0.62, 1.22, 1.15, 0.04, GREEN, radius=False)
    textbox(slide, 0.62, 7.03, 4.0, 0.22, "KHOZO | Proposal for Government of Assam", 7.5, MUTED)
    textbox(slide, 12.0, 7.03, 0.8, 0.22, f"{page:02d}", 7.5, MUTED, align=PP_ALIGN.RIGHT)


def logo(slide, x=0.62, y=0.35, light=False):
    fill_rect(slide, x, y, 0.44, 0.44, WHITE if light else GREEN)
    textbox(slide, x + 0.11, y + 0.07, 0.24, 0.17, "K", 16, DARK_GREEN if light else WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, x + 0.55, y + 0.08, 1.1, 0.25, "KHOZO", 15, WHITE if light else INK, True)


def source(slide, text):
    textbox(slide, 8.4, 6.72, 4.15, 0.22, text, 7.2, MUTED, align=PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, items, size=12.5, color=INK, dot_color=GREEN, row_h=0.48):
    for i, item in enumerate(items):
        yy = y + i * row_h
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.08), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        textbox(slide, x + 0.22, yy, w - 0.22, 0.34, item, size, color)


def card(slide, x, y, w, h, title, body, accent=GREEN, title_size=14, body_size=10.5):
    fill_rect(slide, x, y, w, h, WHITE)
    fill_rect(slide, x, y, 0.08, h, accent, radius=False)
    textbox(slide, x + 0.25, y + 0.16, w - 0.45, 0.28, title, title_size, INK, True)
    textbox(slide, x + 0.25, y + 0.55, w - 0.45, h - 0.65, body, body_size, MUTED)


def stat(slide, x, y, w, label, value, sub, accent=GREEN):
    fill_rect(slide, x, y, w, 1.1, WHITE)
    textbox(slide, x + 0.18, y + 0.13, w - 0.3, 0.22, label.upper(), 8, MUTED, True)
    textbox(slide, x + 0.18, y + 0.4, w - 0.3, 0.34, value, 21, accent, True)
    textbox(slide, x + 0.18, y + 0.82, w - 0.3, 0.2, sub, 7.5, MUTED)


def process(slide, x, y, steps):
    step_w = 2.35
    gap = 0.34
    colors = [BLUE, GREEN, AMBER, DARK_GREEN]
    for i, (title, body) in enumerate(steps):
        xx = x + i * (step_w + gap)
        fill_rect(slide, xx, y, step_w, 1.48, WHITE)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(xx + 0.16), Inches(y + 0.16), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors[i]
        badge.line.fill.background()
        textbox(slide, xx + 0.235, y + 0.215, 0.18, 0.13, str(i + 1), 8, WHITE, True, PP_ALIGN.CENTER)
        textbox(slide, xx + 0.58, y + 0.15, step_w - 0.72, 0.26, title, 12.3, INK, True)
        textbox(slide, xx + 0.18, y + 0.6, step_w - 0.32, 0.62, body, 9.2, MUTED)
        if i < len(steps) - 1:
            textbox(slide, xx + step_w + 0.06, y + 0.55, 0.22, 0.22, ">", 16, GREEN, True, PP_ALIGN.CENTER)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide, 0, 0, 13.333, 7.5, DARK_GREEN, radius=False)
    fill_rect(slide, 7.65, 0, 5.683, 7.5, INK, radius=False)
    for x, y, s, c in [(8.35, 0.9, 2.65, RGBColor(139, 195, 74)), (9.45, 2.25, 2.5, GREEN), (8.4, 4.1, 2.2, RGBColor(53, 96, 16))]:
        oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
        oval.fill.solid()
        oval.fill.fore_color.rgb = c
        oval.line.fill.background()
    logo(slide, light=True)
    textbox(slide, 0.7, 1.42, 5.9, 0.32, "Government Proposal Deck", 14, LIGHT_GREEN, True)
    textbox(slide, 0.62, 1.95, 6.6, 1.3, "Project Khozo for Government of Assam", 37, WHITE, True)
    textbox(slide, 0.68, 3.48, 6.25, 0.75, "A web and mobile platform to help Assam Police, WCD/SCPS, NGOs and citizens find missing children, women and persons using photo reporting and face recognition.", 16.5, RGBColor(235, 245, 230))
    textbox(slide, 0.68, 6.58, 6.2, 0.24, "Based on the Khozo proposal document and original Khozo pitch deck", 8.5, RGBColor(218, 235, 210))
    source(slide, f"{SRC_PDF}, slide 1; {SRC_DOCX}, Overview")


def executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Executive Summary", "Proposal", 2)
    textbox(slide, 0.78, 1.58, 5.9, 0.9, "Children are the nation's assets. Khozo proposes a one-click public reporting and agency response platform to help missing children return home faster.", 20, INK, True)
    bullet_list(slide, 0.95, 3.0, 5.8, [
        "Citizens, NGOs and police upload photos through the Khozo portal or mobile app.",
        "A centralized secure repository stores missing-person photographs and related records.",
        "The matching service compares uploaded images and notifies nearby police stations.",
        "Parents receive SMS or Khozo app pop-up alerts after authorized match confirmation.",
    ], 12.8)
    card(slide, 7.25, 1.62, 4.75, 2.55, "Proposal for Assam", "Adapt the original Government/Police proposal for Assam Government, Assam Police, WCD/SCPS, district child-protection teams, NGOs and citizens.", BLUE, 16, 12)
    card(slide, 7.25, 4.45, 4.75, 1.25, "Core promise", '"Your 1 click - A missing child can return home."', GREEN, 17, 13)
    source(slide, f"{SRC_DOCX}, Overview; {SRC_PDF}, slide 1")


def objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Objectives", "Need and intent", 3)
    items = [
        ("Help parents and police", "Find missing children while reducing manual human effort."),
        ("Create an online service", "Use face detection, face recognition and supporting information to find missing children."),
        ("Develop mobile access", "Provide a mobile application for field reporting by the public and agencies."),
        ("Enable public participation", "Allow citizens to click pictures of suspicious/missing children and upload them to Khozo server."),
    ]
    for i, (t, b) in enumerate(items):
        card(slide, 0.78 + (i % 2) * 5.95, 1.55 + (i // 2) * 2.0, 5.35, 1.42, t, b, [GREEN, BLUE, AMBER, DARK_GREEN][i], 15, 11.5)
    textbox(slide, 1.02, 6.05, 11.1, 0.34, "For Assam, these objectives translate into faster reporting, faster verification and a measurable government response loop.", 15, INK, True, PP_ALIGN.CENTER)
    source(slide, f"{SRC_PDF}, slide 2")


def scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Project Scope", "What Khozo will build", 4)
    textbox(slide, 0.82, 1.55, 5.7, 0.62, "The project scope is to build a web and mobile app that helps the Police department, NGOs and parents find missing children, women or persons using face recognition.", 17, INK, True)
    bullet_list(slide, 0.95, 2.72, 5.8, [
        "Web portal for government, police, NGOs and public users",
        "Android app on Play Store and iOS app on App Store as free downloads",
        "Secure centralized cloud database for photographs and missing-person records",
        "Fast and accurate photo or video search against missing-person face images",
    ], 12.8)
    card(slide, 7.3, 1.55, 4.85, 1.4, "Assam adaptation", "Replace the proposal's Goa Government / Goa Police rollout language with Government of Assam / Assam Police ownership.", BLUE, 15, 11)
    card(slide, 7.3, 3.25, 4.85, 1.4, "Primary users", "Police, parents, citizens, NGOs and government administrators work through role-based access layers.", GREEN, 15, 11)
    source(slide, f"{SRC_DOCX}, Project Scope and Implementation Plan")


def workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Workflow", "How the response works", 5)
    process(slide, 0.7, 1.65, [
        ("Register", "Police register FIR; parents can also register missing-child details."),
        ("Upload", "Public uploads a child picture using the Khozo mobile app or website."),
        ("Match", "Khozo service loads the picture to server and searches the missing-person repository."),
        ("Alert", "SMS or Khozo app pop-up alert is sent to parents and nearby police station."),
    ])
    card(slide, 0.95, 4.35, 3.6, 1.35, "Photo/video search", "The algorithm identifies a missing child, woman or person in photos or video using a private repository.", BLUE)
    card(slide, 4.9, 4.35, 3.6, 1.35, "Central database", "Police photographs and missing-person records are stored securely for matching and review.", GREEN)
    card(slide, 8.85, 4.35, 3.6, 1.35, "Agency notification", "Matched leads are routed to nearby police stations for action and parent notification.", AMBER)
    source(slide, f"{SRC_PDF}, slide 3; {SRC_DOCX}, Workflow")


def accessibility_layers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Service Accessibility Layers", "Users and rights", 6)
    layers = [
        ("Super Admin", "Assam Government / Police Commissioner: all rights, overall view, drill-down, coordination messages", 0.85, 1.58, 11.5, GREEN),
        ("Admin", "Assistant Commissioner / district command: view usage under police stations in jurisdiction, drill down to users", 1.25, 2.62, 10.7, BLUE),
        ("User", "Police station: register missing complaints/FIRs, view registered complaints, send SMS to parents", 1.72, 3.66, 9.75, AMBER),
        ("Public", "Citizens, parents and NGOs: click and upload child picture using Khozo mobile app or portal", 2.2, 4.7, 8.8, DARK_GREEN),
    ]
    for title, body, x, y, w, color in layers:
        fill_rect(slide, x, y, w, 0.72, color)
        textbox(slide, x + 0.2, y + 0.13, 1.65, 0.24, title, 13, WHITE, True)
        textbox(slide, x + 2.05, y + 0.13, w - 2.25, 0.24, body, 10.2, WHITE)
    textbox(slide, 1.0, 6.05, 11.2, 0.34, "The hierarchy mirrors the original Khozo service accessibility model while assigning Assam-specific ownership.", 14.5, INK, True, PP_ALIGN.CENTER)
    source(slide, f"{SRC_PDF}, slide 5; {SRC_DOCX}, Access rules")


def law_enforcement_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Outcome for Law Enforcement Agencies", "Dashboards", 7)
    card(slide, 0.82, 1.52, 3.55, 1.55, "Super Admin dashboard", "Overall Khozo usage view, drill-down to Admin/User levels, email or WhatsApp coordination with teams.", GREEN)
    card(slide, 4.75, 1.52, 3.55, 1.55, "Admin dashboard", "Jurisdiction view across police stations and ability to drill down to user level.", BLUE)
    card(slide, 8.68, 1.52, 3.55, 1.55, "Police station user", "Register missing-child or missing-person complaints and view Khozo usage under station jurisdiction.", AMBER)
    stat(slide, 1.15, 3.8, 2.6, "Primary action", "FIR", "registration and tracking", GREEN)
    stat(slide, 4.15, 3.8, 2.6, "Response", "SMS", "to parents after match", BLUE)
    stat(slide, 7.15, 3.8, 2.6, "Coordination", "Email/WA", "from dashboards", AMBER)
    stat(slide, 10.15, 3.8, 2.2, "Visibility", "Drill-down", "state to station", DARK_GREEN)
    textbox(slide, 0.95, 5.72, 11.4, 0.35, "For Assam Police, the same dashboard model can support state, district and police-station command accountability.", 14.5, INK, True, PP_ALIGN.CENTER)
    source(slide, f"{SRC_DOCX}, Outcome of usage of Khozo app for Law enforcement agencies")


def implementation_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Implementation Plan", "Roles and responsibilities", 8)
    card(slide, 0.82, 1.5, 3.6, 1.75, "Aegis / Khozo", "Host the Khozo portal and release Android/iOS apps as free downloads. Update the app as data security rules evolve.", GREEN)
    card(slide, 4.85, 1.5, 3.6, 1.75, "Government of Assam", "Provide formal ownership, designate nodal departments, approve rollout and align welfare/police channels.", BLUE)
    card(slide, 8.88, 1.5, 3.6, 1.75, "Assam Police", "Provide missing-person photographs/data, register FIRs, verify matches and coordinate local response.", AMBER)
    card(slide, 2.05, 4.0, 3.9, 1.45, "NGOs and citizens", "Use the portal/app to upload sightings and help society participate in finding missing children.", DARK_GREEN)
    card(slide, 6.9, 4.0, 3.9, 1.45, "Awareness channels", "Publicize through welfare schemes, law enforcement, media, TV, radio, internet and public campaigns.", GREEN)
    source(slide, f"{SRC_DOCX}, Implementation Plan and Stakeholders")


def beneficiary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Beneficiaries", "Public value", 9)
    textbox(slide, 0.86, 1.55, 5.55, 0.78, "If a child or person is missing, it affects the child and the parent deeply. Khozo gives families hope by creating a faster path from sighting to official action.", 18, INK, True)
    bullet_list(slide, 1.0, 2.82, 5.7, [
        "Parents can register missing-child details and receive credentials for updates.",
        "Citizens get a simple way to report sightings instead of only informing police in person.",
        "NGOs can support society by using the same portal and mobile app.",
        "Police receive structured records, photographs and match leads in one workflow.",
    ], 12.8)
    card(slide, 7.25, 1.72, 4.85, 1.35, "Family impact", "A missing-child platform is not only a technology system; it is a hope and trust system for families.", GREEN, 15, 11.5)
    card(slide, 7.25, 3.45, 4.85, 1.35, "Government impact", "A statewide rollout can strengthen citizen participation, police response and welfare coordination.", BLUE, 15, 11.5)
    source(slide, f"{SRC_DOCX}, Beneficiaries; {SRC_PDF}, slide 2")


def assam_pilot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Proposed Assam Pilot", "Rollout proposal", 10)
    phases = [
        ("Phase 1", "Approval and setup", "Nominate state nodal owners, confirm pilot districts, define FIR/data-sharing and escalation SOPs."),
        ("Phase 2", "Portal and app pilot", "Configure Assam roles, onboard police stations and NGOs, run citizen sighting and parent registration flows."),
        ("Phase 3", "Awareness campaign", "Use welfare channels, police outreach, media, radio, internet and public endorsements to drive adoption."),
        ("Phase 4", "Review and scale", "Measure cases registered, sightings uploaded, match review time, parent alerts and district adoption."),
    ]
    for i, (phase, title, body) in enumerate(phases):
        y = 1.45 + i * 1.18
        textbox(slide, 0.86, y + 0.12, 0.8, 0.24, phase, 10.5, GREEN, True)
        card(slide, 1.95, y, 9.75, 0.78, title, body, [GREEN, BLUE, AMBER, DARK_GREEN][i], 13, 9.5)
    source(slide, f"{SRC_DOCX}, Implementation Plan; adapted for Assam Government")


def proposal_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Proposal Ask from Government of Assam", "Decision", 11)
    textbox(slide, 0.88, 1.52, 5.65, 0.74, "Approve Khozo as a pilot public-good platform for missing children, women and persons across selected Assam districts.", 20, INK, True)
    card(slide, 0.95, 2.85, 3.55, 1.28, "1. Approve pilot", "Authorize a time-bound pilot with Assam Police and WCD/SCPS participation.", GREEN)
    card(slide, 4.9, 2.85, 3.55, 1.28, "2. Provide data support", "Enable police-approved missing-person photographs and FIR/workflow alignment.", BLUE)
    card(slide, 8.85, 2.85, 3.55, 1.28, "3. Enable public adoption", "Promote the free app through welfare, police, media and NGO channels.", AMBER)
    textbox(slide, 1.0, 5.55, 11.1, 0.44, "Government support can turn Khozo from a technology platform into a statewide missing-person response network.", 18, DARK_GREEN, True, PP_ALIGN.CENTER)
    source(slide, f"{SRC_DOCX}, Implementation Plan and Project Scope")


def references_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "References Used", "Appendix", 12)
    bullet_list(slide, 0.92, 1.55, 11.2, [
        "Khozo (1).pdf, slide 1: initiative by Aegis School of Data Science, Mumbai; tagline and khozo.org reference.",
        "Khozo (1).pdf, slide 2: objectives for parents, police, online face detection/recognition service, mobile app and public photo uploads.",
        "Khozo (1).pdf, slide 3: workflow diagram covering FIR registration, public upload, server loading and SMS/app alert to parents.",
        "Khozo (1).pdf, slide 5: service accessibility layers: Super Admin, Admin and Public.",
        "khozo_proposal_v1.0.docx: background, project scope, workflow, stakeholder roles, implementation plan, law-enforcement outcomes and beneficiaries.",
        "Adaptation note: the DOCX refers to Goa Government and Goa Police; this deck adapts the same proposal structure for Government of Assam and Assam Police.",
    ], 11.2, row_h=0.64)
    textbox(slide, 0.92, 6.08, 11.2, 0.34, "The deck is a government proposal deck, not a verbatim reproduction of the source documents.", 12.5, MUTED)
    source(slide, f"{SRC_PDF}; {SRC_DOCX}")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in [
        title_slide,
        executive_summary,
        objectives,
        scope_slide,
        workflow_slide,
        accessibility_layers,
        law_enforcement_slide,
        implementation_plan,
        beneficiary_slide,
        assam_pilot,
        proposal_ask,
        references_slide,
    ]:
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
