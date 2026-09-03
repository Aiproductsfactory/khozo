"""Build the Khozo proposal deck for the Department of Women and Child
Development, Government of Maharashtra.

Re-angles the Assam police-facing deck for a child-welfare audience: Mission
Vatsalya alignment, the CWC/DCPU/CCI/JJB stakeholder ecosystem, Juvenile
Justice compliance, child privacy and departmental MIS.

    python scripts/create_maharashtra_wcd_deck.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "khozo-maharashtra-wcd-deck.pptx"
SCREEN_DIR = ROOT / "docs" / "screenshots"

GREEN = RGBColor(103, 164, 38)
DARK_GREEN = RGBColor(78, 128, 25)
LIGHT_GREEN = RGBColor(234, 244, 221)
INK = RGBColor(28, 37, 48)
MUTED = RGBColor(86, 99, 115)
BLUE = RGBColor(37, 99, 235)
AMBER = RGBColor(217, 119, 6)
TEAL = RGBColor(13, 148, 136)
WHITE = RGBColor(255, 255, 255)

DEPT = "Department of Women and Child Development, Government of Maharashtra"
SRC_AUDIT = "Source: docs/system-audit.md"

# The addressee. Named on the title slide in the courtesy convention a proposal
# submitted to a Minister's office follows — this deck is addressed TO the
# office, and nothing here should be read as that office endorsing Khozo.
# The portrait is optional: drop the file in and the title slide uses it,
# leave it out and the slide keeps its original treatment. A deck that will not
# build because an image is missing is a deck that fails the night before it is
# needed.
MINISTER_NAME = "Hon'ble Smt. Aditi Sunil Tatkare"
MINISTER_TITLE = "Cabinet Minister, Department of Women and Child Development"
MINISTER_PHOTO = ROOT / "docs" / "assets" / "minister-wcd-maharashtra.jpg"


def textbox(slide, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def fill_rect(slide, x, y, w, h, color, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def header(slide, title, section, page):
    textbox(slide, 0.62, 0.32, 5.2, 0.25, section.upper(), 8.5, GREEN, True)
    textbox(slide, 0.58, 0.58, 10.2, 0.55, title, 24, INK, True)
    fill_rect(slide, 0.62, 1.22, 1.15, 0.04, GREEN, radius=False)
    textbox(slide, 0.62, 7.03, 6.0, 0.22, "KHOZO | Proposal to WCD, Government of Maharashtra", 7.5, MUTED)
    textbox(slide, 12.0, 7.03, 0.8, 0.22, f"{page:02d}", 7.5, MUTED, align=PP_ALIGN.RIGHT)


def logo(slide, x=0.62, y=0.35, light=False):
    fill_rect(slide, x, y, 0.44, 0.44, WHITE if light else GREEN)
    textbox(slide, x + 0.11, y + 0.07, 0.24, 0.17, "K", 16, DARK_GREEN if light else WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, x + 0.55, y + 0.08, 1.1, 0.25, "KHOZO", 15, WHITE if light else INK, True)


def source(slide, text):
    textbox(slide, 8.4, 6.72, 4.15, 0.22, text, 7.2, MUTED, align=PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, items, size=12.5, color=INK, dot_color=GREEN, row_h=0.48):
    for i, item in enumerate(items):
        yy = y + i * row_h
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.08), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        textbox(slide, x + 0.22, yy, w - 0.22, 0.34, item, size, color)


def card(slide, x, y, w, h, title, body, accent=GREEN, title_size=14, body_size=10.5):
    fill_rect(slide, x, y, w, h, WHITE)
    fill_rect(slide, x, y, 0.08, h, accent, radius=False)
    textbox(slide, x + 0.25, y + 0.16, w - 0.45, 0.28, title, title_size, INK, True)
    textbox(slide, x + 0.25, y + 0.55, w - 0.45, h - 0.65, body, body_size, MUTED)


def stat(slide, x, y, w, label, value, sub, accent=GREEN):
    fill_rect(slide, x, y, w, 1.1, WHITE)
    textbox(slide, x + 0.18, y + 0.13, w - 0.3, 0.22, label.upper(), 8, MUTED, True)
    textbox(slide, x + 0.18, y + 0.4, w - 0.3, 0.34, value, 20, accent, True)
    textbox(slide, x + 0.18, y + 0.82, w - 0.3, 0.2, sub, 7.5, MUTED)


def chip(slide, x, y, w, h, label, body, accent):
    fill_rect(slide, x, y, w, h, WHITE)
    fill_rect(slide, x, y, w, 0.06, accent, radius=False)
    textbox(slide, x + 0.16, y + 0.16, w - 0.3, 0.24, label, 11.5, INK, True)
    textbox(slide, x + 0.16, y + 0.46, w - 0.3, h - 0.55, body, 8.4, MUTED)


def process(slide, x, y, steps, step_w=2.28, gap=0.26):
    colors = [BLUE, GREEN, AMBER, TEAL, DARK_GREEN]
    for i, (title, body) in enumerate(steps):
        xx = x + i * (step_w + gap)
        fill_rect(slide, xx, y, step_w, 1.55, WHITE)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(xx + 0.16), Inches(y + 0.16), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors[i % len(colors)]
        badge.line.fill.background()
        textbox(slide, xx + 0.235, y + 0.215, 0.18, 0.13, str(i + 1), 8, WHITE, True, PP_ALIGN.CENTER)
        textbox(slide, xx + 0.58, y + 0.15, step_w - 0.72, 0.26, title, 11.8, INK, True)
        textbox(slide, xx + 0.18, y + 0.62, step_w - 0.32, 0.78, body, 8.8, MUTED)
        if i < len(steps) - 1:
            textbox(slide, xx + step_w + 0.02, y + 0.58, 0.22, 0.22, ">", 15, GREEN, True, PP_ALIGN.CENTER)


# --------------------------------------------------------------------------- slides


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide, 0, 0, 13.333, 7.5, DARK_GREEN, radius=False)
    fill_rect(slide, 7.65, 0, 5.683, 7.5, INK, radius=False)
    if MINISTER_PHOTO.exists():
        # The portrait replaces the decorative ovals rather than sitting on top
        # of them: a Minister's photograph competing with three coloured circles
        # reads as clip art, which is not how a proposal to a Department should
        # present the person it is addressed to.
        box_w, box_h = 2.95, 3.6
        with Image.open(MINISTER_PHOTO) as im:
            iw, ih = im.size
        scale = min(box_w / iw, box_h / ih)
        w, h = iw * scale, ih * scale
        px = 7.65 + (5.683 - w) / 2
        py = 1.9
        slide.shapes.add_picture(str(MINISTER_PHOTO), Inches(px), Inches(py), width=Inches(w), height=Inches(h))
        textbox(slide, 7.95, py + h + 0.24, 5.1, 0.32, MINISTER_NAME, 14, WHITE, True, PP_ALIGN.CENTER)
        textbox(slide, 7.95, py + h + 0.6, 5.1, 0.72, f"{MINISTER_TITLE}\nGovernment of Maharashtra", 10.5, RGBColor(218, 235, 210), False, PP_ALIGN.CENTER)
    else:
        for x, y, s, c in [(8.35, 0.9, 2.65, RGBColor(139, 195, 74)), (9.45, 2.25, 2.5, GREEN), (8.4, 4.1, 2.2, RGBColor(53, 96, 16))]:
            oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
            oval.fill.solid()
            oval.fill.fore_color.rgb = c
            oval.line.fill.background()
    logo(slide, light=True)
    textbox(slide, 0.7, 1.36, 6.4, 0.32, "Proposal Deck | Child Protection Technology", 13.5, LIGHT_GREEN, True)
    textbox(slide, 0.62, 1.85, 6.7, 1.45, "Project Khozo for the Department of Women and Child Development", 33, WHITE, True)
    textbox(slide, 0.68, 3.62, 6.35, 0.9, "A Mission Vatsalya-aligned platform to trace, verify, protect and restore missing children — connecting citizens, NGOs, Police, CWC, DCPU, CCIs and JJBs in one accountable, privacy-safe workflow.", 14.5, RGBColor(235, 245, 230))
    textbox(slide, 0.68, 4.95, 6.35, 0.3, "Government of Maharashtra", 15, WHITE, True)
    textbox(slide, 0.68, 6.42, 6.3, 0.5, "An initiative by Aegis School of Data Science & AI, Mumbai — AI for Social Good\nSubmitted for the kind consideration of " + MINISTER_NAME + ", " + MINISTER_TITLE, 9, RGBColor(218, 235, 210))


def why_this_matters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "The Problem We Are Trying to Solve", "Context", 2)
    textbox(slide, 0.8, 1.5, 6.0, 0.95, "The gap is rarely goodwill. It is the time and traceability between a child going missing, a citizen spotting that child, and the right officer or Committee being able to act on it.", 18, INK, True)
    bullet_list(slide, 0.95, 2.72, 6.0, [
        "A citizen who spots a child often has no safe, structured way to report it.",
        "Sightings and case records sit in separate registers, files and departments.",
        "Welfare follow-up after recovery — CWC production, CCI care, restoration — is tracked manually.",
        "Families have no reliable way to check progress without visiting an office.",
        "Child identity and photographs need protection at every step of this chain.",
    ], 11.8, row_h=0.52)
    card(slide, 7.35, 1.55, 4.75, 1.55, "The consequence", "Delay in the first hours, and loss of a single verified trail from sighting to restoration that the Department can supervise.", AMBER, 15, 11)
    card(slide, 7.35, 3.3, 4.75, 1.5, "What is needed", "Not a replacement for Mission Vatsalya, but a fast intake and coordination layer that feeds it and is answerable to it.", GREEN, 15, 11)
    card(slide, 7.35, 5.0, 4.75, 1.35, "State context", "[Insert current Maharashtra missing-children figures from NCRB / State Crime Records Bureau before presentation.]", BLUE, 15, 10.5)


def executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Executive Summary", "Proposal", 3)
    textbox(slide, 0.78, 1.5, 6.0, 0.85, "Khozo is a working platform that turns a citizen's photograph into an accountable, jurisdiction-routed child-protection action.", 19, INK, True)
    bullet_list(slide, 0.95, 2.6, 6.0, [
        "Citizens, parents and NGOs report a missing child or a sighting from a phone.",
        "Reports are verified by a scoped official before becoming a formal case.",
        "Sightings are ranked and routed to the correct district desk automatically.",
        "Unmatched sightings escalate to Childline 1098 / CWC rather than being lost.",
        "Welfare, restoration and closure are tracked to a supervised outcome.",
    ], 11.8, row_h=0.52)
    card(slide, 7.35, 1.55, 4.75, 1.75, "Built for this Department", "Native roles for CWC, DCPU, SJPU, AHTU, CCI, JJB, DLSA, SARA, SAA and the State Nodal Officer — not a police tool with welfare bolted on.", BLUE, 15, 11)
    card(slide, 7.35, 3.5, 4.75, 1.6, "Complements, not duplicates", "Cases carry TrackChild, KhoyaPaya, CCTNS/FIR, NCRB and GHAR identifiers so existing Government systems stay authoritative.", TEAL, 15, 11)
    card(slide, 7.35, 5.3, 4.75, 1.05, "Core promise", '"Your 1 click — a missing child can return home."', GREEN, 15, 12)


def govt_alignment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Grounded in the Government's Own Framework", "Design basis", 4)
    textbox(slide, 0.8, 1.46, 11.4, 0.4, "Khozo's workflows, roles and records were designed after a structured study of the Government's existing child-protection systems and statutes.", 15, INK, True)
    refs = [
        ("Mission Vatsalya", "Stakeholder structure, grievance and feedback routing, and welfare service categories follow the scheme's operating model.", GREEN),
        ("TrackChild", "Case registers, matching, district search, MIS reporting and identity linkage patterns.", BLUE),
        ("KhoyaPaya", "Citizen-facing missing/found reporting that feeds institutional review instead of exposing child identity publicly.", AMBER),
        ("GHAR (NCPCR)", "Restoration and repatriation tracking, institutional handover, transfer and support needs.", TEAL),
        ("CARINGS / CARA", "Adoption governance separating SARA and SAA responsibilities and legally-free declarations.", DARK_GREEN),
        ("JJ Act, 2015", "Production before CWC/JJB within statutory timelines, Social Investigation Reports and Individual Care Plans.", BLUE),
    ]
    for i, (t, b, c) in enumerate(refs):
        chip(slide, 0.8 + (i % 3) * 3.92, 2.05 + (i // 3) * 2.05, 3.6, 1.82, t, b, c)
    textbox(slide, 0.8, 6.28, 11.4, 0.35, "Khozo is positioned as an intake and coordination layer that strengthens these systems — it does not seek to replace any of them.", 13, DARK_GREEN, True, PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "The Khozo Workflow", "How a case moves", 5)
    process(slide, 0.72, 1.55, [
        ("Report", "Parent, citizen or NGO reports a missing child or uploads a sighting photograph, with consent recorded."),
        ("Verify", "A scoped police / CWC / DCPU officer verifies the intake before it becomes a formal case with an FIR or GD number."),
        ("Match", "The sighting is ranked against registered cases and routed to the correct district desk."),
        ("Refer", "Unmatched or at-risk sightings escalate to Childline 1098, CWC or DCPU rather than closing silently."),
        ("Restore", "Production, care, welfare referral and restoration are tracked to a formal, recorded closure."),
    ])
    card(slide, 0.8, 3.55, 3.72, 1.42, "Nothing is lost", "Every intake either becomes a case, a referral or a recorded disposition. There is no silent dead end.", GREEN)
    card(slide, 4.82, 3.55, 3.72, 1.42, "Nothing is unowned", "Cases carry a jurisdiction, an assigned desk and a transfer history when investigation moves.", BLUE)
    card(slide, 8.84, 3.55, 3.72, 1.42, "Nothing is untraceable", "Every action writes a tamper-evident audit event that the Department can later verify.", AMBER)
    textbox(slide, 0.8, 5.32, 11.7, 0.75, "Citizens can check a receipt or FIR reference at any time and see safe progress — without any child's identity, photograph, address or guardian contact being exposed.", 14.5, INK, True, PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def stakeholder_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "The Child-Protection Ecosystem, Built In", "Stakeholder roles", 6)
    textbox(slide, 0.8, 1.44, 11.4, 0.36, "Each role below is implemented with its own scoped dashboard, permissions, jurisdiction and audit trail.", 14.5, INK, True)
    roles = [
        ("CWC", "Child Welfare Committee — production, orders, care decisions", GREEN),
        ("DCPU", "District Child Protection Unit — district coordination and review", GREEN),
        ("CCI", "Child Care Institution — admission, care plan and service register", TEAL),
        ("JJB", "Juvenile Justice Board — proceedings, hearings and directions", TEAL),
        ("SJPU", "Special Juvenile Police Unit — child-sensitive policing", BLUE),
        ("AHTU", "Anti Human Trafficking Unit — trafficking-risk cases", BLUE),
        ("DLSA", "District Legal Services Authority — legal aid and compensation", AMBER),
        ("SARA / SAA", "Adoption governance and CARINGS-aligned records", AMBER),
        ("RPF / Railways", "Railway and station-based rescue coordination", DARK_GREEN),
        ("DCRB / SCRB", "District and State Crime Records Bureau reporting", DARK_GREEN),
        ("State Nodal Officer", "State-level escalation, oversight and MIS", GREEN),
        ("Citizens / NGOs", "Verified public reporting and sighting participation", BLUE),
    ]
    for i, (t, b, c) in enumerate(roles):
        chip(slide, 0.8 + (i % 4) * 2.95, 1.94 + (i // 4) * 1.55, 2.72, 1.35, t, b, c)
    textbox(slide, 0.8, 6.62, 11.4, 0.3, "Reunification confirmation remains restricted to Police / SJPU and command roles. Welfare roles hold welfare powers — authority is never blurred.", 11.5, MUTED, align=PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def statutory_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Statutory Compliance as a Product Feature", "Juvenile Justice workflow", 7)
    textbox(slide, 0.8, 1.46, 11.4, 0.36, "Compliance obligations the Department already carries are tracked, with due dates and overdue alerts.", 14.5, INK, True)
    items = [
        ("24-hour production", "Rescue time, production time, authority, order number and outcome are recorded, with the statutory deadline computed and breaches flagged.", GREEN),
        ("SIR and ICP", "Social Investigation Reports and Individual Care Plans are structured records with assessor, risk level and review dates.", BLUE),
        ("CCI care register", "Admission type, institution, care plan, service categories and next review date, with a district-level institutional register.", TEAL),
        ("JJB proceedings", "Proceeding type, board, case number, order date, next hearing and directions.", AMBER),
        ("Restoration and repatriation", "Route, handover authority, escort, documents, funding and follow-up — tracked through to closure.", DARK_GREEN),
        ("Formal closure", "Restored to family, CCI transfer, adoption or aftercare, repatriation, or traced pending reunification.", GREEN),
    ]
    for i, (t, b, c) in enumerate(items):
        chip(slide, 0.8 + (i % 3) * 3.92, 1.98 + (i // 3) * 2.1, 3.6, 1.88, t, b, c)
    textbox(slide, 0.8, 6.28, 11.4, 0.34, "A follow-up queue surfaces overdue production, due reviews and stale cases so supervision is proactive rather than retrospective.", 13, DARK_GREEN, True, PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def privacy_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Child Privacy Is a Design Constraint", "Safeguards", 8)
    textbox(slide, 0.8, 1.5, 6.05, 0.85, "A missing-child platform handles the most sensitive data the State holds. Khozo was built so that openness to citizens never becomes exposure of a child.", 17, INK, True)
    bullet_list(slide, 0.95, 2.62, 6.05, [
        "Public screens redact guardian contacts, Aadhaar, exact addresses and photographs.",
        "Photographs require recorded consent, stated purpose and a retention period.",
        "Photograph access is authenticated and jurisdiction-checked — never a public URL.",
        "Citizens may mark their identity confidential; review queues show masked details.",
        "A privacy review queue governs retention, extension and anonymisation.",
        "Anonymisation requires an approval or order reference and is fully audited.",
    ], 11.4, row_h=0.5)
    card(slide, 7.4, 1.58, 4.7, 1.5, "Audit metadata is redacted too", "Audit events store flags, counts and text lengths — not names, phone numbers, ID numbers or note contents.", GREEN, 14.5, 10.5)
    card(slide, 7.4, 3.28, 4.7, 1.5, "Tamper-evident record", "Audit events form a hash chain. Exports carry a SHA-256 digest and HMAC signature that can be verified after download.", BLUE, 14.5, 10.5)
    card(slide, 7.4, 4.98, 4.7, 1.45, "Abuse protection", "Public intake is rate-limited, OTP-verified, and abuse signals feed a command review queue with hashed identities.", AMBER, 14.5, 10.5)
    source(slide, SRC_AUDIT)


def oversight_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Oversight and MIS for the Department", "Command view", 9)
    card(slide, 0.8, 1.5, 3.72, 1.6, "State command view", "Statewide caseload, reunification, pending intake, district breakdown and workflow activity for the Minister's office and State Nodal Officer.", GREEN, 14.5, 10.5)
    card(slide, 4.82, 1.5, 3.72, 1.6, "District view", "DCPU and district desks see only their jurisdiction — cases, sightings, grievances, reviews and coverage gaps.", BLUE, 14.5, 10.5)
    card(slide, 8.84, 1.5, 3.72, 1.6, "Signed MIS report", "A scoped MIS artifact with a verifiable signature, generated on demand and audited each time.", TEAL, 14.5, 10.5)
    stat(slide, 0.95, 3.4, 2.65, "Coverage", "Gap report", "missing district desks", GREEN)
    stat(slide, 3.9, 3.4, 2.65, "Supervision", "Follow-up", "overdue and stale work", AMBER)
    stat(slide, 6.85, 3.4, 2.65, "Evidence", "Signed export", "audit and MIS artifacts", BLUE)
    stat(slide, 9.8, 3.4, 2.65, "Coordination", "Network alert", "scoped desk messaging", TEAL)
    textbox(slide, 0.85, 4.95, 11.5, 0.8, "The coverage report tells the Department which district desks — CWC, DCPU, SJPU, CCI, JJB, RPF — are not yet onboarded in a State or district, so gaps are visible before a pilot begins rather than after it fails.", 14.5, INK, True, PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def citizen_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Citizen Participation, Safely Channelled", "Public interface", 10)
    card(slide, 0.8, 1.5, 3.72, 1.75, "Emergency routing", "Citizens are routed by context to Childline 1098, ERSS 112 or Railway 139, paired with the correct local SJPU, CWC, DCPU or RPF desk.", GREEN, 14.5, 10.5)
    card(slide, 4.82, 1.5, 3.72, 1.75, "Public service directory", "A redacted directory of local child-protection response points — without exposing official accounts or internal user data.", BLUE, 14.5, 10.5)
    card(slide, 8.84, 1.5, 3.72, 1.75, "Privacy-safe bulletins", "Authorised officers publish missing-child alerts with guardian contacts, Aadhaar, addresses and protected photographs redacted.", AMBER, 14.5, 10.5)
    card(slide, 0.8, 3.55, 3.72, 1.75, "Grievance channel", "A Mission Vatsalya-style grievance and feedback route with a reference number and safe public status tracking.", TEAL, 14.5, 10.5)
    card(slide, 4.82, 3.55, 3.72, 1.75, "Status without exposure", "Families check a receipt, FIR or TrackChild reference and see progress — never another child's identity.", DARK_GREEN, 14.5, 10.5)
    card(slide, 8.84, 3.55, 3.72, 1.75, "Works in weak network", "Field reports queue in encrypted on-device storage and sync when connectivity returns — important for tribal and rural blocks.", GREEN, 14.5, 10.5)
    textbox(slide, 0.85, 5.62, 11.5, 0.6, "Public registration is OTP-verified and privileged roles cannot be self-assigned — official accounts are provisioned only by authorised command users.", 14, INK, True, PP_ALIGN.CENTER)
    source(slide, SRC_AUDIT)


def sovereign_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Face Recognition Under Indian Sovereignty", "Data sovereignty", 11)
    textbox(slide, 0.8, 1.46, 6.15, 0.95, "Khozo proposes Sovereign AI Aarakshak for production face recognition, so children's biometric data stays under Indian infrastructure and governance.", 16, INK, True)
    bullet_list(slide, 0.95, 2.62, 6.15, [
        "Indian child-protection data should remain in India, under Indian institutions.",
        "No sharing of children's facial-recognition data with foreign companies.",
        "Matching is isolated behind a single service boundary in the codebase.",
        "A human officer always reviews and confirms a match — the system never decides alone.",
        "Every match decision is recorded as signed, audit-ready evidence.",
    ], 11.6, row_h=0.53)
    # Was: "the current build uses a non-biometric demonstration scorer".
    # That stopped being true when the fallback scorer was removed and the
    # two-engine rule shipped; Aarakshak and AWS Rekognition both run in the
    # deployed system. Telling a Department we are less capable than we are is
    # as much a misstatement as the reverse, and this slide is the one the
    # Department will hold us to.
    card(slide, 7.45, 1.6, 4.65, 1.55, "Honest position today", "Face matching runs on Aarakshak, an Indian provider, with AWS Rekognition (Mumbai region) as an independent second opinion. A candidate is shown to an officer only where both engines agree; anything uncorroborated is withheld, not shown with a lower number.", AMBER, 14.5, 9.5)
    card(slide, 7.45, 3.35, 4.65, 1.5, "Integration model", "Matching lives behind one module, so an authorised provider replaces the demo scorer without changing intake, dashboards or safeguards.", BLUE, 14.5, 10.5)
    card(slide, 7.45, 5.05, 4.65, 1.35, "Departmental assurance", "Face recognition assists officers on child-protection cases only — it is not a surveillance capability.", GREEN, 14.5, 10.5)
    source(slide, "Sources: aarakshak.com; server/src/match.js")


def screenshot_slide(prs, title, filename, caption, page, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, title, "Working portal — actual screenshot", page)
    img_path = SCREEN_DIR / filename
    box_w, box_h = 8.6, 5.25
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = iw * scale, ih * scale
    slide.shapes.add_picture(str(img_path), Inches(0.75), Inches(1.5), width=Inches(w), height=Inches(h))
    card(slide, 9.6, 1.55, 2.9, 2.5, "What this shows", caption, accent, 13, 10)
    textbox(slide, 9.72, 4.3, 2.6, 0.6, f"docs/screenshots/{filename}", 8.5, MUTED)
    source(slide, "Captured from the running Khozo portal")


def status_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Where the Platform Stands Today", "Honest status", 18)
    textbox(slide, 0.8, 1.48, 11.4, 0.4, "We wish to be candid with the Department about maturity, so that any pilot is planned on accurate ground.", 15, INK, True)
    card(slide, 0.8, 2.02, 5.72, 2.35, "Working today", "End-to-end citizen intake, verification, sighting review and referral. All stakeholder roles with scoped dashboards. Statutory and welfare records. Privacy, retention and anonymisation controls. Tamper-evident audit with signed exports. Grievances, bulletins, public directory and offline field capture.", GREEN, 15, 10.5)
    card(slide, 6.8, 2.02, 5.72, 2.35, "To be completed before live use", "Production database in place of the demonstration store. Authorised face recognition connected via Aarakshak. Live SMS gateway for OTP and parent alerts. Government-grade hosting, security audit and key management. Full automated test coverage and operational support.", AMBER, 15, 10.5)
    textbox(slide, 0.85, 4.7, 11.5, 0.85, "Khozo is a fully demonstrable, end-to-end prototype developed by faculty and student teams at Aegis School of Data Science & AI. It is offered to the Government in the spirit of public service. We seek the Department's direction on a supervised district pilot to validate it against real operational conditions before any wider deployment.", 14, INK, True, PP_ALIGN.CENTER)
    textbox(slide, 0.85, 5.75, 11.5, 0.4, "No live child data has been used in development. All demonstration records are synthetic.", 13, DARK_GREEN, True, PP_ALIGN.CENTER)
    source(slide, "Source: docs/system-audit.md, docs/session-handoff.md")


def pilot_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "Proposed Maharashtra Pilot", "Rollout", 19)
    phases = [
        ("Phase 1", "Approval and governance", "Department nominates the State Nodal Officer and pilot districts; SOPs agreed for verification, referral, data sharing and escalation.", "4–6 weeks"),
        ("Phase 2", "Onboarding and hardening", "Provision CWC, DCPU, CCI, JJB, SJPU and police desks; production database, hosting, security review and SMS gateway put in place.", "6–8 weeks"),
        ("Phase 3", "Supervised district pilot", "Live use in the selected districts with citizen and NGO participation, weekly review with DCPU and the State Nodal Officer.", "3–4 months"),
        ("Phase 4", "Evaluation and scale decision", "Measure verification time, referral closure, production compliance and restoration outcomes; Department decides on wider rollout.", "4–6 weeks"),
    ]
    for i, (phase, title, body, dur) in enumerate(phases):
        y = 1.5 + i * 1.28
        textbox(slide, 0.82, y + 0.14, 0.9, 0.24, phase, 10.5, GREEN, True)
        card(slide, 1.85, y, 8.55, 0.98, title, body, [GREEN, BLUE, AMBER, TEAL][i], 13, 9.5)
        textbox(slide, 10.65, y + 0.34, 1.85, 0.3, dur, 11.5, MUTED, True, PP_ALIGN.CENTER)
    textbox(slide, 0.85, 6.7, 11.5, 0.32, "Indicative timeline, offered for the Department's revision. Aegis will bear platform development and support costs for the pilot period.", 12, DARK_GREEN, True, PP_ALIGN.CENTER)


def ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "What We Seek from the Department", "Our request", 20)
    textbox(slide, 0.82, 1.5, 11.4, 0.45, "We are not seeking funds. We seek direction, supervision and the authority to be useful.", 18, INK, True, PP_ALIGN.CENTER)
    card(slide, 0.8, 2.25, 3.72, 1.9, "1. Guidance on a pilot", "Nominate one or two districts and a nodal officer under whose supervision Khozo can be validated in real conditions.", GREEN, 15, 11)
    card(slide, 4.82, 2.25, 3.72, 1.9, "2. Workflow validation", "Allow CWC, DCPU and CCI officers to review whether the workflows and records match how the Department actually works.", BLUE, 15, 11)
    card(slide, 8.84, 2.25, 3.72, 1.9, "3. Integration direction", "Advise on alignment with Mission Vatsalya and TrackChild, and on authorisation for face recognition.", AMBER, 15, 11)
    textbox(slide, 0.85, 4.55, 11.5, 0.9, "With the Department's guidance, a student-built public-good platform can become a supervised State capability — one that shortens the time between a citizen's photograph and a child's safe return home.", 17, DARK_GREEN, True, PP_ALIGN.CENTER)
    textbox(slide, 0.85, 5.75, 11.5, 0.35, "We would be glad to demonstrate the working system to the Department at any date convenient to the Hon'ble Minister's office.", 13.5, MUTED, align=PP_ALIGN.CENTER)


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide, 0, 0, 13.333, 7.5, DARK_GREEN, radius=False)
    fill_rect(slide, 7.7, 0, 5.633, 7.5, INK, radius=False)
    textbox(slide, 0.78, 0.78, 1.2, 0.28, "KHOZO", 15, WHITE, True)
    textbox(slide, 0.78, 1.72, 6.5, 1.35, "Let every safe sighting become an accountable response.", 33, WHITE, True)
    textbox(slide, 0.82, 3.42, 6.1, 0.75, "A proposal to the Department of Women and Child Development, Government of Maharashtra, for a supervised, privacy-safe, Mission Vatsalya-aligned child-protection pilot.", 15, LIGHT_GREEN)
    card(slide, 0.88, 4.55, 2.92, 1.15, "Next step", "Guidance on a supervised district pilot and a nodal officer.", WHITE, 14, 10)
    card(slide, 4.1, 4.55, 2.92, 1.15, "Offered by", "Aegis School of Data Science & AI, Mumbai — AI for Social Good.", WHITE, 14, 10)
    textbox(slide, 0.88, 6.1, 6.1, 0.55, "[Presenter name] | [Designation] | [Mobile] | [Email]", 11, LIGHT_GREEN)
    textbox(slide, 8.35, 2.05, 3.55, 0.5, "Your 1 click", 28, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, 8.35, 2.74, 3.55, 0.85, "A missing child can return home", 24, LIGHT_GREEN, True, PP_ALIGN.CENTER)
    textbox(slide, 8.15, 5.72, 3.95, 0.28, "An initiative by Aegis School of Data Science & AI", 9.5, MUTED, align=PP_ALIGN.CENTER)


SHOTS = [
    ("Command Dashboard", "overview.png", "Role-scoped caseload, reunification, pending matches, readiness and follow-up alerts for command and district users.", GREEN),
    ("Case and Intake Register", "cases-firs.png", "Operational case list with missing-child intake, verification and FIR/GD registration workflows.", BLUE),
    ("CCI Institutional Register", "cci-register.png", "Admissions, active children, admission-type split and review-due status across Child Care Institutions.", TEAL),
    ("Privacy and Retention Review", "privacy-review.png", "Retention status, review, extension and approved anonymisation for case and sighting records.", AMBER),
    ("Signed MIS Report", "mis-report.png", "Scoped MIS artifact with district breakdowns, compliance counts and a verifiable signature.", DARK_GREEN),
    ("Audit Trail", "audit-log.png", "Hash-chained audit events with integrity verification and signed export for evidence handling.", BLUE),
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for fn in [
        title_slide,
        why_this_matters,
        executive_summary,
        govt_alignment,
        workflow_slide,
        stakeholder_slide,
        statutory_slide,
        privacy_slide,
        oversight_slide,
        citizen_slide,
        sovereign_slide,
    ]:
        fn(prs)

    for i, (title, filename, caption, accent) in enumerate(SHOTS, start=12):
        screenshot_slide(prs, title, filename, caption, i, accent)

    status_slide(prs)
    pilot_slide(prs)
    ask_slide(prs)
    closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"{OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
