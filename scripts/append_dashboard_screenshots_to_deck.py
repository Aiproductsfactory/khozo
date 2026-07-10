from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from scripts.create_assam_govt_deck import (
    OUT, build, header, source, textbox, card, bullet_list, fill_rect,
    GREEN, BLUE, AMBER, DARK_GREEN, LIGHT_GREEN, INK, MUTED, WHITE,
)

SCREEN_DIR = ROOT / 'docs' / 'screenshots'

FEATURES = [
    'Implemented: role-scoped overview dashboard with charts, pilot readiness and follow-up alerts.',
    'Implemented: Cases & FIRs, child/FIR registration, public sightings and match review.',
    'Implemented: stakeholder network, account provisioning, coverage gaps and network alerts.',
    'Implemented: CCI register, privacy review, audit log, signed MIS report and public abuse queue.',
    'Implemented: grievances and public-facing receipt/status workflows alongside demo login roles.',
    'Deck gap fixed: added real dashboard screenshots so the proposal shows the product already built.',
]

SHOT_SLIDES = [
    ('Dashboard Overview', 'overview.png', 'Role-scoped command dashboard with caseload, reunification, pending matches, readiness, follow-ups and charts.'),
    ('Cases, FIRs And Match Review', 'cases-firs.png', 'Operational case list paired with missing-child/FIR workflows and reviewable sightings/matches.'),
    ('MIS, Audit And Privacy Controls', 'mis-report.png', 'Signed MIS report plus audit/privacy surfaces for traceability, retention and government reporting.'),
    ('Stakeholder Network And Coverage', 'network.png', 'Provision stakeholder accounts, send coordination alerts and inspect state/district coverage gaps.'),
    ('Public Abuse And Grievance Queues', 'public-abuse.png', 'Command review queues for public intake abuse signals, OTP/rate-limit events and safe follow-up.'),
]


def add_sovereign_ai_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 'Sovereign AI Facial Recognition With Aarakshak', 'Data sovereignty', 13)
    textbox(
        slide,
        0.82,
        1.46,
        6.25,
        0.72,
        'Khozo will use Sovereign AI Aarakshak for production facial recognition, so sensitive biometric and child-protection data remains under India-aligned infrastructure and governance.',
        18,
        INK,
        True,
    )
    bullet_list(slide, 0.95, 2.58, 6.6, [
        'Aarakshak positions itself as sovereign verified presence for India\'s institutions.',
        'Production recognition can be connected behind Khozo\'s existing match-provider boundary.',
        'Indian child-protection and identity data should remain in India and not be shared with foreign companies.',
        'Human review gates, signed evidence packets and audit-ready decisions fit government workflows.',
        'The current Khozo demo scorer is non-biometric; Aarakshak is the proposed production recognition provider.',
    ], 11.9, row_h=0.53)
    card(slide, 7.7, 1.62, 4.25, 1.32, 'Government assurance', 'Sovereign deployment keeps control with Indian institutions while allowing Assam Police and WCD/SCPS to use face recognition responsibly.', GREEN, 14, 10.5)
    card(slide, 7.7, 3.22, 4.25, 1.32, 'Integration model', 'Khozo already isolates matching in server/src/match.js, so Aarakshak can replace the demo-local scorer without changing public intake or dashboards.', BLUE, 14, 10.5)
    card(slide, 7.7, 4.82, 4.25, 1.0, 'Message for Assam', 'India data remains in India. No sharing of facial-recognition data with foreign companies.', AMBER, 14, 10.5)
    source(slide, 'Sources: aarakshak.com; server/src/match.js')


def add_feature_coverage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, 'Built Portal Feature Coverage', 'Product verification', 14)
    textbox(slide, 0.82, 1.48, 6.2, 0.48, 'Checked against the actual Khozo portal routes and captured authenticated super-admin dashboard screenshots.', 17, INK, True)
    bullet_list(slide, 0.95, 2.25, 6.4, FEATURES, 11.6, row_h=0.52)
    card(slide, 7.65, 1.58, 4.35, 1.25, 'Current deck fit', 'The main proposal deck represents the original Khozo scope well: portal/app, face-match workflow, FIRs, dashboards, access layers and public participation.', GREEN, 14, 10.5)
    card(slide, 7.65, 3.08, 4.35, 1.45, 'What screenshots add', 'The appendix now proves the newer built modules: MIS, audit, privacy, abuse, CCI, grievances and network coverage.', BLUE, 14, 10.5)
    card(slide, 7.65, 4.78, 4.35, 1.05, 'Screenshot account', 'Captured from superadmin@khozo.org against the local built portal.', AMBER, 14, 10.5)
    source(slide, 'Sources: web/src/dashboard/routes.js and docs/screenshots/*.png')


def add_screenshot_slide(prs, title, filename, caption, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, title, 'Actual portal screenshot', page)
    img = SCREEN_DIR / filename
    slide.shapes.add_picture(str(img), Inches(0.74), Inches(1.48), width=Inches(8.65))
    card(slide, 9.75, 1.55, 2.75, 2.2, 'What this proves', caption, GREEN if page % 2 else BLUE, 13, 10.2)
    textbox(slide, 9.86, 4.1, 2.42, 0.72, f'File: docs/screenshots/{filename}', 9.5, MUTED)
    source(slide, 'Captured from actual Khozo portal running locally')


def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide, 0, 0, 13.333, 7.5, DARK_GREEN, radius=False)
    fill_rect(slide, 7.7, 0, 5.633, 7.5, INK, radius=False)
    textbox(slide, 0.78, 0.78, 1.2, 0.28, 'KHOZO', 15, WHITE, True)
    textbox(slide, 0.78, 1.72, 6.5, 1.25, 'Let every safe sighting become an accountable response.', 34, WHITE, True)
    textbox(slide, 0.82, 3.35, 6.0, 0.65, 'Proposal to partner with Government of Assam, Assam Police and WCD/SCPS for a sovereign, India-hosted missing-child response pilot.', 16, LIGHT_GREEN)
    card(slide, 0.9, 4.55, 2.9, 1.05, 'Next step', 'Approve pilot and nominate state nodal owners.', WHITE, 14, 10.5)
    card(slide, 4.15, 4.55, 2.9, 1.05, 'Technology', 'Khozo portal with Aarakshak sovereign facial recognition.', WHITE, 14, 10.5)
    textbox(slide, 8.35, 2.05, 3.55, 0.5, 'Your 1 click', 28, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, 8.35, 2.74, 3.55, 0.78, 'A missing child can return home', 24, LIGHT_GREEN, True, PP_ALIGN.CENTER)
    textbox(slide, 8.15, 5.72, 3.95, 0.28, 'An initiative by Aegis School of Data Science', 9.5, MUTED, align=PP_ALIGN.CENTER)


def main():
    build()
    prs = Presentation(str(OUT))
    add_sovereign_ai_slide(prs)
    add_feature_coverage(prs)
    for offset, (title, file, caption) in enumerate(SHOT_SLIDES, start=15):
        add_screenshot_slide(prs, title, file, caption, offset)
    add_closing_slide(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == '__main__':
    main()

